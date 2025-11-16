"""
Script untuk mengkonversi HTML documentation ke PowerPoint (PPTX)
Menggunakan python-pptx library
"""

import os
import re
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background gradient effect (using shape)
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    
    # Background rectangle
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)  # Dark blue-gray
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(203, 213, 225)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Version info
    version_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    version_frame = version_box.text_frame
    version_frame.text = "Versi 1.0"
    version_para = version_frame.paragraphs[0]
    version_para.font.size = Pt(14)
    version_para.font.color.rgb = RGBColor(148, 163, 184)
    version_para.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, content_items):
    """Create content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title[:100] if len(title) > 100 else title  # Limit title length
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    title_shape.text_frame.word_wrap = True
    
    # Content
    if len(slide.placeholders) > 1:
        content_shape = slide.placeholders[1]
    else:
        # Fallback: create textbox if placeholder doesn't exist
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        content_shape = slide.shapes.add_textbox(left, top, width, height)
    
    tf = content_shape.text_frame
    tf.word_wrap = True
    tf.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Clean text - remove HTML tags and excessive whitespace
        text = re.sub(r'<[^>]+>', '', str(item))
        text = re.sub(r'\s+', ' ', text).strip()
        
        # Limit text length per bullet
        if len(text) > 150:
            text = text[:147] + "..."
        
        p.text = text
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_after = Pt(10)
        p.space_before = Pt(5)
    
    return slide

def create_section_slide(prs, section_title, section_number):
    """Create section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(37, 99, 235)  # Blue
    bg.line.fill.background()
    
    # Section number
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    num_frame = num_box.text_frame
    num_frame.text = f"Bagian {section_number}"
    num_para = num_frame.paragraphs[0]
    num_para.font.size = Pt(32)
    num_para.font.color.rgb = RGBColor(255, 255, 255)
    num_para.alignment = PP_ALIGN.CENTER
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide

def parse_html_to_slides(html_file):
    """Parse HTML and extract content for slides"""
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    slides_data = []
    
    # Extract cover page info
    cover = soup.find('div', class_='cover-page')
    if cover:
        title = cover.find('h1')
        subtitle = cover.find('div', class_='subtitle')
        slides_data.append({
            'type': 'cover',
            'title': title.text.strip() if title else 'ROC Dashboard',
            'subtitle': subtitle.text.strip() if subtitle else 'Panduan Penggunaan'
        })
    
    # Extract sections from pages
    pages = soup.find_all('div', class_='page')
    
    for page in pages:
        # Get section title (h1)
        h1 = page.find('h1')
        if h1:
            title = h1.get_text(strip=True)
            # Remove emojis and extract section number
            title_clean = re.sub(r'[^\d\s\.]', '', title).strip()
            match = re.match(r'(\d+)', title_clean)
            if match:
                section_num = match.group(1)
                # Get original title without numbers
                section_title = re.sub(r'^\d+\.\s*', '', title).strip()
                # Remove emojis
                section_title = re.sub(r'[^\w\s\-]', '', section_title).strip()
                
                slides_data.append({
                    'type': 'section',
                    'number': section_num,
                    'title': section_title
                })
        
        # Extract h2 subsections
        h2s = page.find_all('h2')
        for h2 in h2s:
            h2_text = h2.get_text(strip=True)
            # Remove emojis and numbers
            h2_clean = re.sub(r'[^\w\s\-]', '', h2_text).strip()
            h2_clean = re.sub(r'^\d+\.\d+\s*', '', h2_clean).strip()
            
            if h2_clean:
                slides_data.append({
                    'type': 'content',
                    'title': h2_clean,
                    'items': []
                })
        
        # Extract feature boxes
        feature_boxes = page.find_all('div', class_='feature-box')
        for box in feature_boxes:
            h3 = box.find('h3')
            if h3:
                box_title = h3.get_text(strip=True)
                # Remove emojis
                box_title = re.sub(r'[^\w\s\-]', '', box_title).strip()
                items = []
                
                # Get list items
                ul = box.find('ul', class_='icon-list')
                if ul:
                    for li in ul.find_all('li'):
                        text = li.get_text(strip=True)
                        # Remove emojis and clean text
                        text = re.sub(r'[^\w\s\-.,;:!?()\[\]{}]', '', text).strip()
                        if text:
                            items.append(text)
                else:
                    # Get paragraph text
                    paragraphs = box.find_all('p')
                    for p in paragraphs:
                        text = p.get_text(strip=True)
                        if text:
                            # Remove emojis
                            text = re.sub(r'[^\w\s\-.,;:!?()\[\]{}]', '', text).strip()
                            if text and len(text) > 10:  # Only add substantial text
                                items.append(text)
                
                if items:
                    # Split into multiple slides if too many items
                    max_items = 6
                    for i in range(0, len(items), max_items):
                        slide_items = items[i:i+max_items]
                        slide_title = box_title
                        if i > 0:
                            slide_title += f" (Lanjutan {i//max_items + 1})"
                        
                        slides_data.append({
                            'type': 'content',
                            'title': slide_title,
                            'items': slide_items
                        })
        
        # Extract step boxes
        step_boxes = page.find_all('div', class_='step-box')
        for box in step_boxes:
            h3 = box.find('h3')
            if h3:
                box_title = h3.get_text(strip=True)
                # Remove emojis
                box_title = re.sub(r'[^\w\s\-]', '', box_title).strip()
                items = []
                
                for p in box.find_all('p'):
                    text = p.get_text(strip=True)
                    if text:
                        # Remove step number circles and emojis
                        text = re.sub(r'^\d+', '', text).strip()
                        text = re.sub(r'[^\w\s\-.,;:!?()\[\]{}]', '', text).strip()
                        if text and len(text) > 5:
                            items.append(text)
                
                if items:
                    # Split into multiple slides if too many items
                    max_items = 6
                    for i in range(0, len(items), max_items):
                        slide_items = items[i:i+max_items]
                        slide_title = box_title
                        if i > 0:
                            slide_title += f" (Lanjutan {i//max_items + 1})"
                        
                        slides_data.append({
                            'type': 'content',
                            'title': slide_title,
                            'items': slide_items
                        })
    
    return slides_data

def create_powerpoint(html_file, output_file):
    """Main function to create PowerPoint from HTML"""
    print(f"Membaca file HTML: {html_file}")
    
    if not os.path.exists(html_file):
        print(f"❌ File {html_file} tidak ditemukan!")
        return False
    
    try:
        # Parse HTML
        print("Memproses konten HTML...")
        slides_data = parse_html_to_slides(html_file)
        
        if not slides_data:
            print("⚠️ Tidak ada konten yang ditemukan dalam HTML")
            return False
        
        # Create presentation
        print("Membuat PowerPoint presentation...")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Create slides
        slide_count = 0
        for data in slides_data:
            if data['type'] == 'cover':
                create_title_slide(prs, data['title'], data['subtitle'])
                slide_count += 1
            elif data['type'] == 'section':
                create_section_slide(prs, data['title'], data['number'])
                slide_count += 1
            elif data['type'] == 'content':
                create_content_slide(prs, data['title'], data['items'])
                slide_count += 1
        
        # Save presentation
        print(f"Menyimpan PowerPoint: {output_file}")
        prs.save(output_file)
        
        print(f"✅ Berhasil! PowerPoint dibuat dengan {slide_count} slide")
        print(f"📄 File tersimpan di: {output_file}")
        return True
        
    except ImportError:
        print("❌ Library yang diperlukan tidak terinstall!")
        print("Install dengan: pip install python-pptx beautifulsoup4")
        return False
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    html_file = "PANDUAN_PENGGUNAAN_ROC_DASHBOARD.html"
    output_file = "PANDUAN_PENGGUNAAN_ROC_DASHBOARD.pptx"
    
    print("=" * 60)
    print("Konversi HTML ke PowerPoint - ROC Dashboard")
    print("=" * 60)
    print()
    
    if create_powerpoint(html_file, output_file):
        print()
        print("=" * 60)
        print("Selesai! Buka file PowerPoint untuk melihat hasilnya.")
        print("=" * 60)
    else:
        print()
        print("=" * 60)
        print("Gagal membuat PowerPoint. Periksa error di atas.")
        print("=" * 60)

if __name__ == "__main__":
    main()

